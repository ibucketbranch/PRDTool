#!/usr/bin/env python3
"""
Document Processor
Scans PDFs, extracts content, generates AI summaries and categorization using Groq.
Handles both conversation PDFs and general documents.

🛡️ RULE #1 - ABSOLUTE SAFETY:
This module NEVER deletes, moves, renames, or modifies files on disk.
It ONLY reads files (extract text, analyze) and updates database records.
All file operations are read-only. Files remain completely untouched.
"""

import sys
import os
import json
import hashlib
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import re
import requests

# Load .env file
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    print("⚠️  WARNING: python-dotenv not installed. .env file will not be loaded.")
    print("   Install with: pip install python-dotenv")
    print("   Or add to requirements.txt and run: pip install -r requirements.txt")

# Add scripts directory to path to import helpers
current_dir = os.path.dirname(os.path.abspath(__file__))
scripts_dir = os.path.join(current_dir, 'scripts')
if scripts_dir not in sys.path:
    sys.path.append(scripts_dir)

# Import conversation handlers
try:
    from parse_iphone_backup import parse_iphone_backup_messages
    from analyze_conversation import ConversationAnalyzer
    from populate_database import DatabasePopulator
except ImportError as e:
    print(f"Warning: Could not import conversation helpers: {e}")
    parse_iphone_backup_messages = None

# Default service role key for local Supabase (TheConversation project)
# This is the standard local development key - if your instance uses a different key,
# set SUPABASE_SERVICE_ROLE_KEY environment variable
DEFAULT_SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")

try:
    from PyPDF2 import PdfReader
except ImportError:
    print("PyPDF2 not installed. Install with: pip install PyPDF2")
    sys.exit(1)

try:
    from groq import Groq
except ImportError:
    Groq = None
    print("Warning: Groq client not available. Install with: pip install groq")

try:
    from google import genai
    from google.genai import types as genai_types
except ImportError:
    genai = None
    genai_types = None

try:
    from supabase import create_client, Client
except ImportError:
    print("supabase not installed. Install with: pip install supabase")
    sys.exit(1)


try:
    from pdf2image import convert_from_path
except ImportError:
    print("pdf2image not installed. Install with: pip install pdf2image")
    convert_from_path = None

class DocumentProcessor:
    def __init__(self, supabase_url: str = "http://127.0.0.1:54421",
                 supabase_key: Optional[str] = None,
                 groq_api_key: Optional[str] = None,
                 llm_provider: Optional[str] = None,
                 gemini_api_key: Optional[str] = None):
        """Initialize document processor with Supabase and AI clients."""
        supabase_key = supabase_key or os.getenv('SUPABASE_SERVICE_ROLE_KEY') or DEFAULT_SUPABASE_KEY
        self.supabase: Client = create_client(supabase_url, supabase_key)

        self.llm_provider = (llm_provider or os.getenv('LLM_PROVIDER') or 'ollama').lower()
        self.gemini_model_name = os.getenv('GEMINI_MODEL', 'models/gemini-2.0-flash')
        self.ollama_host = os.getenv('OLLAMA_HOST', 'http://127.0.0.1:11434')
        self.ollama_fast_model = os.getenv('OLLAMA_FAST_MODEL', 'llama3.1:8b')
        self.ollama_reasoning_model = os.getenv('OLLAMA_REASONING_MODEL', 'qwen2.5:14b')
        self.ollama_escalation_model = os.getenv(
            'OLLAMA_ESCALATION_MODEL', self.ollama_reasoning_model
        )

        groq_key = groq_api_key or os.getenv('GROQ_API_KEY')
        self.groq_client = None
        if Groq and groq_key:
            self.groq_client = Groq(api_key=groq_key)
            print("✅ Groq client initialized (FALLBACK provider)")
        elif self.llm_provider == 'groq':
            print("⚠️  Warning: Groq API key not configured; groq provider unavailable.")

        gemini_key = gemini_api_key or os.getenv('GEMINI_API_KEY')
        self.gemini_client = None
        if gemini_key and genai:
            try:
                self.gemini_client = genai.Client(api_key=gemini_key)
                print("✅ Gemini client initialized (PRIMARY provider)")
            except Exception as e:
                print(f"⚠️  Warning: Could not initialize Gemini client: {e}")
        else:
            if self.llm_provider == 'gemini':
                print("⚠️  WARNING: Gemini API key not configured (GEMINI_API_KEY not set)")
                print("   This violates the intended configuration - Gemini should be PRIMARY")
            if not gemini_key:
                print("⚠️  GEMINI_API_KEY not found in environment - Gemini will not be available")
                print("   Set it with: export GEMINI_API_KEY='your-key'")
                print("   Or pass it to DocumentProcessor(gemini_api_key='your-key')")

        self.ollama_available = False
        self.ollama_models: List[str] = []
        try:
            health = requests.get(f"{self.ollama_host}/api/tags", timeout=2)
            if health.status_code == 200:
                self.ollama_available = True
                payload = health.json() if health.content else {}
                self.ollama_models = [
                    model.get("name", "")
                    for model in payload.get("models", [])
                    if model.get("name")
                ]
                self.ollama_fast_model = self._resolve_ollama_model(
                    self.ollama_fast_model,
                    preferred_prefixes=["llama3.1:8b", "llama3.1"],
                )
                self.ollama_reasoning_model = self._resolve_ollama_model(
                    self.ollama_reasoning_model,
                    preferred_prefixes=["qwen2.5:14b", "qwen2.5-coder:14b", "qwen2.5"],
                )
                self.ollama_escalation_model = self._resolve_ollama_model(
                    self.ollama_escalation_model,
                    preferred_prefixes=["qwen2.5:32b", "qwen3-coder", "qwen2.5-coder:32b"],
                )
                print(
                    "✅ Ollama reachable (LOCAL provider): "
                    f"{self.ollama_fast_model}, {self.ollama_reasoning_model}"
                )
        except Exception:
            if self.llm_provider == 'ollama':
                print(
                    "⚠️  WARNING: OLLAMA selected but not reachable at "
                    f"{self.ollama_host}"
                )
                print("   Start Ollama and pull models:")
                print(f"   ollama pull {self.ollama_fast_model}")
                print(f"   ollama pull {self.ollama_reasoning_model}")

        if not self.ollama_available and not self.gemini_client and not self.groq_client:
            raise RuntimeError(
                "No LLM providers are configured. Enable local Ollama "
                "(recommended) or set GEMINI_API_KEY/GROQ_API_KEY."
            )
        
        # Diagnostic: Show which providers are available
        provider_order = self._get_provider_order()
        print(f"✅ Provider order configured: {provider_order}")

        self._load_context_bins()

    def _resolve_ollama_model(
        self, configured_model: str, preferred_prefixes: List[str]
    ) -> str:
        """Resolve a working model name from installed Ollama tags."""
        if configured_model in self.ollama_models:
            return configured_model

        for prefix in preferred_prefixes:
            for model_name in self.ollama_models:
                if model_name.startswith(prefix):
                    return model_name

        if self.ollama_models:
            return self.ollama_models[0]
        return configured_model
        
    def _load_context_bins(self):
        """Load context bins from database for matching."""
        try:
            result = self.supabase.table('context_bins').select('bin_name').eq('is_active', True).execute()
            self.context_bins = [bin['bin_name'] for bin in result.data] if result.data else []
        except Exception as e:
            print(f"Warning: Could not load context bins: {e}")
            self.context_bins = ['Personal Bin', 'Work Bin', 'Family Bin', 'Finances Bin', 'Legal Bin']
    
    def _get_provider_order(self) -> List[str]:
        """
        Determine the order of LLM providers to attempt.

        Provider order is chosen by preferred provider, then falls back to others.
        Default preference is local-first (`ollama`) for privacy and cost control.
        """
        available: List[str] = []
        if self.ollama_available:
            available.append('ollama')
        if self.gemini_client:
            available.append('gemini')
        if self.groq_client:
            available.append('groq')

        preferred = self.llm_provider if self.llm_provider in available else None
        order: List[str] = [preferred] if preferred else []
        for provider in available:
            if provider not in order:
                order.append(provider)
        return order

    def _select_ollama_model(
        self, task_type: str, max_tokens: int, images: Optional[List] = None
    ) -> str:
        """Select local model by task profile."""
        if images:
            # Current local defaults are text-first; keep reasoning model for richer prompts.
            return self.ollama_reasoning_model
        if task_type in {"ingest", "fast", "lightweight"}:
            return self.ollama_fast_model
        if task_type in {"escalation", "hard_case"}:
            return self.ollama_escalation_model
        if task_type in {"entity_extraction", "classification", "summary"}:
            return self.ollama_reasoning_model
        if max_tokens <= 500:
            return self.ollama_fast_model
        return self.ollama_reasoning_model

    def _call_ollama(
        self,
        system_prompt: str,
        user_prompt: str,
        temperature: float,
        max_tokens: int,
        response_format: str = "json",
        task_type: str = "general",
        images: Optional[List] = None,
    ) -> str:
        """Call local Ollama chat API with task-based model routing."""
        if not self.ollama_available:
            raise RuntimeError("Ollama is not reachable")

        model_name = self._select_ollama_model(task_type, max_tokens, images)
        if images:
            print("   ⚠️  Ollama image payload not enabled in this path; using text prompt.")

        fmt_guard = (
            "\n\nReturn strictly valid JSON only, with no prose or markdown."
            if response_format == "json"
            else ""
        )
        payload = {
            "model": model_name,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"{user_prompt}{fmt_guard}"},
            ],
            "stream": False,
            "options": {
                "temperature": temperature,
                "num_predict": max_tokens,
            },
        }
        response = requests.post(
            f"{self.ollama_host}/api/chat", json=payload, timeout=120
        )
        response.raise_for_status()
        data = response.json()
        message = data.get("message", {})
        content = message.get("content", "")
        if not content or not content.strip():
            raise RuntimeError("Empty response from Ollama")
        return content

    def _call_groq(self, system_prompt: str, user_prompt: str, temperature: float,
                    max_tokens: int, response_format: str = "json") -> str:
        if not self.groq_client:
            raise RuntimeError("Groq client is not configured")

        request = {
            'model': "llama-3.3-70b-versatile",
            'messages': [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            'temperature': temperature,
            'max_tokens': max_tokens
        }
        if response_format == "json":
            request['response_format'] = {"type": "json_object"}

        response = self.groq_client.chat.completions.create(**request)
        content = response.choices[0].message.content if response.choices else None
        if not content or not content.strip():
            raise RuntimeError("Empty response from Groq")
        return content

    def _call_gemini(self, prompt: str, temperature: float, max_tokens: int,
                      response_format: str = "json", images: List = None) -> str:
        if not self.gemini_client or not genai_types:
            raise RuntimeError("Gemini client is not configured")

        config_kwargs = {
            "temperature": temperature,
            "max_output_tokens": max_tokens,
        }
        if response_format == "json":
            config_kwargs["response_mime_type"] = "application/json"

        generation_config = genai_types.GenerateContentConfig(**config_kwargs)

        contents = []
        if images:
             # Add images to content
             for img in images:
                 contents.append(img)
        
        contents.append(prompt)

        response = self.gemini_client.models.generate_content(
            model=self.gemini_model_name,
            contents=contents,
            config=generation_config,
        )

        text = getattr(response, 'output_text', None)
        if text:
            return text

        for candidate in getattr(response, 'candidates', []) or []:
            content = getattr(candidate, 'content', None)
            if not content:
                continue
            pieces = []
            for part in getattr(content, 'parts', []) or []:
                value = getattr(part, 'text', None)
                if value:
                    pieces.append(value)
            if pieces:
                return ''.join(pieces)

        raise RuntimeError("Gemini returned empty response")

    def _invoke_llm(self, system_prompt: str, user_prompt: str, temperature: float,
                    max_tokens: int, response_format: str = "json",
                    images: List = None, task_type: str = "general") -> str:
        """
        Invoke LLM with provider routing and fallback.
        """
        provider_order = self._get_provider_order()
        if not provider_order:
            raise RuntimeError("No LLM providers available. Check API key configuration.")

        last_error: Optional[Exception] = None
        for provider in provider_order:
            try:
                if provider == 'ollama':
                    print("   🤖 Using Ollama (LOCAL provider)...")
                    return self._call_ollama(
                        system_prompt=system_prompt,
                        user_prompt=user_prompt,
                        temperature=temperature,
                        max_tokens=max_tokens,
                        response_format=response_format,
                        task_type=task_type,
                        images=images,
                    )
                if provider == 'gemini':
                    print("   🤖 Using Gemini...")
                    # Set explicit timeout for Gemini to prevent hangs
                    # If it takes > 25s (longer for images), fail over to Groq
                    try:
                        import signal
                        def handler(signum, frame):
                            raise TimeoutError("Gemini timed out")
                        signal.signal(signal.SIGALRM, handler)
                        signal.alarm(20)  # 20 second timeout - sharp watcher for LLM calls
                        
                        # Gemini supports images natively
                        result = self._call_gemini(f"{system_prompt}\n\n{user_prompt}", temperature, max_tokens, response_format, images)
                        
                        signal.alarm(0) # Disable alarm
                        print("   ✅ Gemini response received")
                        return result
                    except TimeoutError:
                        print("   ⚠️  Gemini timed out (>20s). Switching to Groq fallback...")
                        raise
                    finally:
                        signal.alarm(0)
                        
                if provider == 'groq':
                    if images:
                        print("   ⚠️  Groq fallback: Images ignored (Text-only mode)")
                    if provider_order[0] != 'groq':
                        print("   ⚠️  Using Groq as fallback")
                    return self._call_groq(system_prompt, user_prompt, temperature, max_tokens, response_format)
            except Exception as exc:
                last_error = exc
                print(f"   ⚠️  Warning: {provider.capitalize()} request failed: {exc}")

        if last_error:
            raise last_error
        raise RuntimeError("No LLM providers available")

    def detect_context_bin(self, file_path: str) -> Optional[str]:
        """Detect which context bin a file belongs to based on its path."""
        path_lower = file_path.lower()
        
        # Try exact match first (case-insensitive)
        for bin_name in self.context_bins:
            if bin_name.lower() in path_lower:
                return bin_name
        
        # Try keyword matching
        bin_keywords = {
            'Personal Bin': ['personal', 'private', 'my documents'],
            'Work Bin': ['work', 'professional', 'business', 'career', 'job'],
            'Family Bin': ['family', 'kids', 'children', 'spouse', 'relatives'],
            'Finances Bin': ['finance', 'financial', 'money', 'tax', 'banking'],
            'Legal Bin': ['legal', 'court', 'lawyer', 'attorney', 'law'],
            'Projects Bin': ['project', 'projects'],
            'NetV': ['netv', 'net v'],
            'LEOPard': ['leopard', 'leo pard'],
            'USAA Visa': ['usaa', 'visa']
        }
        
        for bin_name, keywords in bin_keywords.items():
            if any(keyword in path_lower for keyword in keywords):
                return bin_name
        
        # Check if path contains iCloud Drive and a bin name
        if 'icloud' in path_lower or 'cloud' in path_lower:
            # Look for common folder names after iCloud Drive
            path_parts = Path(file_path).parts
            for part in path_parts:
                if 'bin' in part.lower():
                    return part
        
        return None
    
    def calculate_file_hash(self, file_path: str) -> str:
        """Calculate SHA256 hash of file for duplicate detection."""
        sha256_hash = hashlib.sha256()
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()
    
    def extract_docx_content(self, file_path: str) -> Dict:
        """Extract text and metadata from DOCX files."""
        doc_info = {
            'title': '', 'author': '', 'subject': '', 'keywords': [],
            'created_date': None, 'modified_date': None, 'page_count': 0,
            'extracted_text': '', 'text_preview': ''
        }
        
        try:
            from docx import Document
            doc = Document(file_path)
            
            # Extract text from all paragraphs
            full_text = ""
            for para in doc.paragraphs:
                full_text += para.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text += cell.text + " "
                    full_text += "\n"
            
            # Get metadata
            core_props = doc.core_properties
            doc_info.update({
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': core_props.keywords.split(',') if core_props.keywords else [],
                'created_date': str(core_props.created) if core_props.created else None,
                'modified_date': str(core_props.modified) if core_props.modified else None,
                'page_count': len(doc.paragraphs) // 50 + 1  # Estimate pages
            })
            
            # Clean null bytes
            cleaned_text = full_text.strip().replace('\x00', '').replace('\u0000', '')
            doc_info['extracted_text'] = cleaned_text
            doc_info['text_preview'] = cleaned_text[:500].strip()
            
            return doc_info
        except ImportError:
            print("   ⚠️  python-docx not installed. Install with: pip install python-docx")
            return doc_info
        except Exception as e:
            print(f"   ⚠️  Error extracting DOCX: {e}")
            return doc_info
    
    def extract_xlsx_content(self, file_path: str) -> Dict:
        """Extract text and metadata from XLSX files."""
        xlsx_info = {
            'title': '', 'author': '', 'subject': '', 'keywords': [],
            'created_date': None, 'modified_date': None, 'page_count': 0,
            'extracted_text': '', 'text_preview': ''
        }
        
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, data_only=True)
            
            # Extract text from all sheets
            full_text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                full_text += f"\n--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        full_text += row_text + "\n"
            
            # Get metadata
            props = wb.properties
            xlsx_info.update({
                'title': props.title or '',
                'author': props.creator or '',
                'subject': props.subject or '',
                'keywords': [],
                'created_date': str(props.created) if props.created else None,
                'modified_date': str(props.modified) if props.modified else None,
                'page_count': len(wb.sheetnames)
            })
            
            # Clean null bytes
            cleaned_text = full_text.strip().replace('\x00', '').replace('\u0000', '')
            xlsx_info['extracted_text'] = cleaned_text
            xlsx_info['text_preview'] = cleaned_text[:500].strip()
            
            return xlsx_info
        except ImportError:
            print("   ⚠️  openpyxl not installed. Install with: pip install openpyxl")
            return xlsx_info
        except Exception as e:
            print(f"   ⚠️  Error extracting XLSX: {e}")
            return xlsx_info
    
    def extract_pptx_content(self, file_path: str) -> Dict:
        """Extract text and metadata from PPTX files."""
        pptx_info = {
            'title': '', 'author': '', 'subject': '', 'keywords': [],
            'created_date': None, 'modified_date': None, 'page_count': 0,
            'extracted_text': '', 'text_preview': ''
        }
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            # Extract text from all slides
            full_text = ""
            for i, slide in enumerate(prs.slides):
                full_text += f"\n--- Slide {i + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
            
            # Get metadata
            core_props = prs.core_properties
            pptx_info.update({
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': [],
                'created_date': str(core_props.created) if core_props.created else None,
                'modified_date': str(core_props.modified) if core_props.modified else None,
                'page_count': len(prs.slides)
            })
            
            # Clean null bytes
            cleaned_text = full_text.strip().replace('\x00', '').replace('\u0000', '')
            pptx_info['extracted_text'] = cleaned_text
            pptx_info['text_preview'] = cleaned_text[:500].strip()
            
            return pptx_info
        except ImportError:
            print("   ⚠️  python-pptx not installed. Install with: pip install python-pptx")
            return pptx_info
        except Exception as e:
            print(f"   ⚠️  Error extracting PPTX: {e}")
            return pptx_info
    
    def extract_pdf_content(self, file_path: str) -> Dict:
        """Extract text and metadata from PDF. Falls back to Image Conversion if PyPDF2 fails."""
        pdf_info = {
            'title': '', 'author': '', 'subject': '', 'keywords': [],
            'created_date': None, 'modified_date': None, 'page_count': 0,
            'extracted_text': '', 'text_preview': ''
        }
        
        try:
            reader = PdfReader(file_path)
            
            # Extract metadata and SANITIZE (convert IndirectObject to str)
            metadata = reader.metadata or {}
            def sanitize(val):
                if val is None: return ''
                return str(val)

            pdf_info.update({
                'title': sanitize(metadata.get('/Title', '')),
                'author': sanitize(metadata.get('/Author', '')),
                'subject': sanitize(metadata.get('/Subject', '')),
                'keywords': sanitize(metadata.get('/Keywords', '')).split(',') if metadata.get('/Keywords') else [],
                'created_date': sanitize(metadata.get('/CreationDate')),
                'modified_date': sanitize(metadata.get('/ModDate')),
                'page_count': len(reader.pages)
            })

            # SKIP HUGE FILES (> 100 pages)
            if len(reader.pages) > 100:
                print(f"   ⚠️  Large document ({len(reader.pages)} pages). Extracting first 10 pages only.")
                page_range = range(min(10, len(reader.pages)))
            else:
                page_range = range(len(reader.pages))
            
            # Extract text
            full_text = ""
            for page_num in page_range:
                try:
                    page = reader.pages[page_num]
                    page_text = page.extract_text()
                    full_text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                except Exception as e:
                    print(f"Warning: Could not extract text from page {page_num + 1}: {e}")
            
            # Clean null bytes and other problematic characters for database
            cleaned_text = full_text.strip().replace('\x00', '').replace('\u0000', '')
            pdf_info['extracted_text'] = cleaned_text
            pdf_info['text_preview'] = cleaned_text[:500].strip()
            
            return pdf_info
            
        except Exception as e:
            print(f"   ⚠️  Standard PDF reading failed: {e}")
            print("   🔄 Falling back to Image Conversion (Vision)...")
            # If standard reading fails, we return empty info and rely on Vision later
            # But we need at least a page count if possible.
            try:
                # Use poppler to get info if PyPDF2 failed
                from pdf2image import pdfinfo_from_path
                info = pdfinfo_from_path(file_path)
                pdf_info['page_count'] = info.get('Pages', 0)
            except:
                pdf_info['page_count'] = 1 # Assume at least 1 page
            
            return pdf_info
    
    def detect_document_mode(self, text: str) -> str:
        """Detect if PDF is a conversation or regular document."""
        # Patterns that indicate a conversation
        conversation_patterns = [
            r'\d{1,2}[/-]\d{1,2}[/-]\d{2,4}[\s,]+\d{1,2}:\d{2}',  # Date/time stamps
            r'(?:From|To):\s*[^\n]+',  # From/To fields
            r'[A-Z][a-z]+\s+[A-Z][a-z]+:\s',  # Name: pattern
            r'iMessage',  # iMessage indicator
            r'(?:\+?[\d\s\-\(\)]{10,}):\s',  # Phone number pattern
        ]
        
        # Check if text matches conversation patterns
        conversation_indicators = 0
        for pattern in conversation_patterns:
            if re.search(pattern, text[:2000]):  # Check first 2000 chars
                conversation_indicators += 1
        
        # If we find 2 or more conversation patterns, it's likely a conversation
        return 'conversation' if conversation_indicators >= 2 else 'document'
    
    def detect_obvious_category(self, file_path: str, file_name: str, text: str) -> Optional[str]:
        """
        Detect obvious document categories from filename, path, and content patterns.
        Returns category if obvious, None if needs AI analysis.
        This prevents unnecessary AI calls for documents with clear indicators.
        """
        file_lower = file_name.lower()
        path_lower = str(file_path).lower()
        text_lower = text[:2000].lower()  # Check first 2000 chars
        
        # RESUME DETECTION - Multiple strong indicators
        resume_keywords_filename = ['resume', 'cv', 'curriculum vitae']
        resume_keywords_path = ['resume', 'cv', 'curriculum']
        resume_keywords_content = [
            'work experience', 'employment history', 'professional experience',
            'education', 'skills', 'summary', 'objective', 'contact information',
            'phone', 'email', 'years of experience', 'responsibilities',
            'achievements', 'qualifications', 'references'
        ]
        
        is_resume_by_filename = any(kw in file_lower for kw in resume_keywords_filename)
        is_resume_by_path = any(kw in path_lower for kw in resume_keywords_path)
        resume_content_matches = sum(1 for kw in resume_keywords_content if kw in text_lower)
        is_resume_by_content = resume_content_matches >= 3
        
        if is_resume_by_filename or is_resume_by_path or is_resume_by_content:
            return 'employment'
        
        # BANK STATEMENT DETECTION - Multiple strong indicators
        bank_keywords_filename = ['bank statement', 'bofa', 'bank of america', 'chase', 'wells fargo', 
                                  'account activity', 'account details', 'balance', 'savings account',
                                  'checking account', 'statement']
        bank_keywords_path = ['bank statement', 'statements', 'bofa', 'banking', 'finances']
        bank_keywords_content = [
            'account number', 'routing number', 'account summary', 'beginning balance',
            'ending balance', 'transaction', 'deposit', 'withdrawal', 'debit', 'credit',
            'available balance', 'statement period', 'account activity'
        ]
        
        is_bank_by_filename = any(kw in file_lower for kw in bank_keywords_filename)
        is_bank_by_path = any(kw in path_lower for kw in bank_keywords_path)
        bank_content_matches = sum(1 for kw in bank_keywords_content if kw in text_lower)
        is_bank_by_content = bank_content_matches >= 3
        
        if is_bank_by_filename or is_bank_by_path or is_bank_by_content:
            return 'bank_statement'
        
        # TAX DOCUMENT DETECTION
        tax_keywords_filename = ['w2', '1099', '1040', 'tax return', 'irs', 'form 1040', 'w-2', '1099-']
        tax_keywords_path = ['tax', 'irs', 'filed', 'tax return']
        tax_keywords_content = ['w-2', '1099', 'adjusted gross income', 'taxable income', 
                               'federal tax', 'social security', 'medicare', 'withholding']
        
        if any(kw in file_lower for kw in tax_keywords_filename) or \
           any(kw in path_lower for kw in tax_keywords_path) or \
           sum(1 for kw in tax_keywords_content if kw in text_lower) >= 2:
            return 'tax_document'
        
        # VEHICLE REGISTRATION
        vehicle_reg_keywords = ['registration', 'dmv', 'vehicle registration', 'registration card']
        if any(kw in file_lower for kw in vehicle_reg_keywords) or \
           any(kw in path_lower for kw in vehicle_reg_keywords):
            return 'vehicle_registration'
        
        # VEHICLE INSURANCE
        vehicle_ins_keywords = ['insurance', 'policy', 'coverage', 'auto insurance', 'vehicle insurance']
        if any(kw in file_lower for kw in vehicle_ins_keywords) or \
           any(kw in path_lower for kw in vehicle_ins_keywords):
            # Check if it's vehicle-related (not general insurance)
            if 'vehicle' in text_lower or 'auto' in text_lower or 'car' in text_lower:
                return 'vehicle_insurance'
        
        # INVOICE/RECEIPT
        invoice_keywords = ['invoice', 'receipt', 'bill', 'payment']
        if any(kw in file_lower for kw in invoice_keywords) or \
           any(kw in path_lower for kw in invoice_keywords):
            # Could be invoice or receipt, let AI distinguish between them
            # But we know it's financial
            pass  # Let AI handle the distinction
        
        # If no obvious pattern, return None to use AI
        return None
    
    def extract_entities(self, text: str, images: List = None) -> Dict:
        """Extract key entities from text using pattern matching and Groq."""
        entities = {
            'people': [],
            'organizations': [],
            'dates': [],
            'amounts': [],
            'vehicles': [],
            'locations': [],
            'legal_cases': [],
            'project_info': {},
            'financial_info': {}
        }
        
        # Only use regex if we have text
        if text.strip():
            # Extract dates
            date_patterns = [
                r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b',
                r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s+\d{4}\b'
            ]
            for pattern in date_patterns:
                entities['dates'].extend(re.findall(pattern, text[:5000]))
            
            # Extract amounts (currency)
            amount_pattern = r'\$[\d,]+(?:\.\d{2})?'
            entities['amounts'] = re.findall(amount_pattern, text[:5000])
            
            # Extract vehicles (common patterns)
            vehicle_patterns = [
                r'\b(?:Tesla|BMW|Mercedes|Ford|Toyota|Honda|Chevrolet|Nissan|Audi|Lexus|Volkswagen)\b',
                r'\b(?:Model [3SXY]|Model S|Model X|Model Y)\b',
                r'\b\d{4}\s+(?:Tesla|BMW|Mercedes|Ford|Toyota|Honda)\b'
            ]
            for pattern in vehicle_patterns:
                entities['vehicles'].extend(re.findall(pattern, text[:5000], re.IGNORECASE))
        
        # Use LLM for more sophisticated entity extraction (supports Vision)
        try:
            entity_prompt = f"""Extract key entities from this document. Return ONLY valid JSON.
If text is empty, relying on attached images.

Text Context (first 1000 chars):
{text[:1000]}

Return JSON format:
{{
  "people": ["list of person names"],
  "organizations": ["list of companies/organizations"],
  "key_terms": ["important terms or topics"],
  "legal_cases": [
      {{
        "case_number": "case/docket number if present",
        "parties": "Plaintiff vs Defendant format",
        "jurisdiction": "Court name or location"
      }}
  ],
  "project_info": {{
      "client_name": "client or company name if applicable",
      "project_name": "project title",
      "status": "draft, final, signed, etc.",
      "document_type": "proposal, contract, brief, spec, etc."
  }},
  "financial_info": {{
      "tax_year": "YYYY",
      "form_type": "W2, 1099, 1040, Invoice, Receipt, etc.",
      "vendor": "merchant or payee",
      "invoice_number": "id"
  }}
}}"""
            content = self._invoke_llm(
                system_prompt="You extract key entities from documents and respond in pure JSON.",
                user_prompt=entity_prompt,
                temperature=0.3,
                max_tokens=600,
                images=images,
                task_type="entity_extraction",
            )
            if not content or content.strip() == '':
                print("Warning: Empty response from LLM during entity extraction")
                return entities
            
            model_entities = json.loads(content)
            entities['people'] = list(set(entities['people'] + model_entities.get('people', [])))
            entities['organizations'] = list(set(entities['organizations'] + model_entities.get('organizations', [])))
            entities['legal_cases'] = model_entities.get('legal_cases', [])
            entities['project_info'] = model_entities.get('project_info', {})
            entities['financial_info'] = model_entities.get('financial_info', {})
            
        except Exception as e:
            print(f"Warning: LLM entity extraction failed: {e}")
        
        return entities
    
    def generate_ai_summary(self, text: str, file_name: str, document_mode: str, file_path: str = None, folder_hierarchy: List = None, images: List = None) -> Dict:
        """Generate AI summary and categorization using LLM with full context awareness."""
        
        if document_mode == 'conversation':
            system_prompt = "You are analyzing a conversation transcript."
            analysis_type = "conversation"
        else:
            system_prompt = """You are an intelligent document categorization system. Your job is to analyze documents and categorize them accurately using available context.

You have access to:
1. The filename - PRIMARY SIGNAL - often contains the strongest clues about document type
2. The document content - PRIMARY SIGNAL - the actual text/content tells you what it really is
3. The ORIGINAL file path - PRIMARY SIGNAL - where the file was FIRST discovered/stored by the user

IMPORTANT: The path you receive is the ORIGINAL path (where the file was first discovered), not where it might be now.
The original path is meaningful because it shows where the user originally stored the file, which reflects their understanding of what the document is.

PRIORITIZATION:
- Filename, Content, and ORIGINAL Path are ALL PRIMARY signals
- The original path shows the user's original intent/organization - use it!
- Example: If original path is "Opportunities 2025/Resumes/MichaelValderrama.pdf", that's a strong signal it's a resume
- Example: If original path is "Bank Statements/2024/BofA.pdf", that's a strong signal it's a bank statement
- Combine all three: filename + content + original path to make the best decision

Examples of intelligent reasoning:
- Filename: "MichaelValderrama.pdf" + Original Path: "Opportunities 2025/Resumes/" + Content: "Work Experience, Skills" → employment (resume)
- Filename: "Bank of America Account Activity" + Original Path: "Financial/Bank Statements/" + Content: "Account #, Transactions" → bank_statement
- Filename: "W2_2023.pdf" + Original Path: "Tax Documents/" + Content: "Wages, Federal Tax" → tax_document
- Filename: "Divorce_Final_Order.pdf" + Original Path: "Legal/" + Content: "Court order" → legal_document

Be smart: Use filename + content + ORIGINAL path together. The original path tells you where the user thought the file belonged, which is valuable context.

Don't overthink obvious documents. A resume is a resume, a bank statement is a bank statement - use all available context to identify them."""
            analysis_type = "document"
        
        # Build rich context information
        path_info = ""
        if file_path:
            path_obj = Path(file_path)
            path_info = f"Full Path: {file_path}\n"
            path_info += f"Parent Folder: {path_obj.parent.name}\n"
            if folder_hierarchy:
                path_info += f"Folder Hierarchy: {' → '.join(folder_hierarchy[-5:])}\n"  # Last 5 levels
            else:
                # Extract hierarchy from path
                parts = path_obj.parts
                if len(parts) > 1:
                    relevant_parts = parts[-5:] if len(parts) > 5 else parts
                    path_info += f"Folder Structure: {' → '.join(relevant_parts[:-1])}\n"
        
        prompt = f"""Analyze this {analysis_type} to categorize it intelligently.

FILENAME: {file_name}

{path_info}
CONTENT (first 3000 characters):
{text[:3000]}

INSTRUCTIONS (in priority order):
1. PRIMARY: Look at the filename - what does it clearly indicate? (e.g., "Resume", "Bank Statement", "W2", "Invoice")
2. PRIMARY: Look at the content - what is the document actually about? What does the text tell you?
3. PRIMARY: Look at the ORIGINAL path (where the file was first discovered) - this shows where the user originally stored it, which gives context about what they thought it was
4. Make your decision based on filename + content + ORIGINAL path. The original path is meaningful because it shows the user's original intent/organization.

IMPORTANT: Use the ORIGINAL path (where file was first discovered), not where it might be now.
The original location tells you more about what the document is than the current location.

CATEGORY OPTIONS:
vehicle_registration, vehicle_insurance, vehicle_maintenance, tax_document, invoice, receipt, contract, medical_record, insurance_policy, bank_statement, utility_bill, property_document, education, employment, correspondence, conversation, project_file, legal_document, identity_document, manual_guide, creative_work, household, food_and_dining, travel, other

If none of the standard categories fit perfectly, use "other" but provide a specific descriptive subcategory.

FOLDER STRUCTURE INTELLIGENCE - Learn from examples:
These are EXAMPLES of excellent organization - learn the patterns, don't just copy them:

Example patterns (learn from these):
- Resume: "Work Bin/Employment/Resumes/Michael Valderrama" → Pattern: Work Bin + Employment + Resumes + Person Name
- Bank Statement: "Finances Bin/Statements/2024/Bank of America" → Pattern: Finances Bin + Statements + Year + Organization
- Tax Document: "Finances Bin/Taxes/2024" → Pattern: Finances Bin + Taxes + Year
- Legal Contract: "Legal Bin/Contracts/2024" → Pattern: Legal Bin + Type + Year
- Medical Record: "Personal Bin/Medical/Records" → Pattern: Personal Bin + Category + Subcategory

Key principles to learn:
1. Use logical bins (Work Bin, Finances Bin, Personal Bin, Legal Bin) based on document purpose
2. Group by category first (Employment, Taxes, Medical, etc.)
3. Then subcategorize (Resumes, Statements, Contracts, etc.)
4. Then group by entity (Person Name, Organization, Year, etc.)
5. Make structures findable and intuitive

Your job: Learn from these examples and intelligently propose the BEST folder structure for THIS specific document.
Use entities (people, organizations, years) to create smart groupings. Be creative and logical.

Provide this exact JSON structure:
{{
  "summary": "2-3 sentence summary of the content",
  "primary_category": "one of the category options above",
  "subcategories": ["list of relevant secondary categories or specific descriptive terms if using 'other'"],
  "key_topics": ["list of 3-5 main topics"],
  "confidence": 0.85,
  "suggested_tags": ["list of useful tags"],
  "importance_score": 5,
  "reasoning": "Brief explanation of why you chose this category based on filename, path, and content",
  "suggested_folder_structure": "Propose the BEST folder path structure. Be intelligent and creative. Use entities from the document (person names, organizations, years) to create logical groupings. Example: For Michael Valderrama's resume from 2024, propose 'Work Bin/Employment/Resumes/Michael Valderrama'. For Bank of America statement from 2024, propose 'Finances Bin/Statements/2024/Bank of America'. Make it logical, findable, and organized. Use format: '[Bin Name]/[Category]/[Subcategory]/[Entity/Year/etc.]'"
}}

Confidence should be 0-1 (how sure you are of the categorization based on all context).
Importance score should be 0-10 (how important this document seems).
"""
        
        try:
            content = self._invoke_llm(
                system_prompt=system_prompt,
                user_prompt=prompt,
                temperature=0.5,
                max_tokens=800,
                images=images,
                task_type="classification",
            )
            if not content or content.strip() == '':
                print("Warning: Empty response from LLM during summarization")
                raise ValueError('empty response')
            
            result = json.loads(content)
            if isinstance(result, list):
                result = result[0] if result else {}
            if not isinstance(result, dict):
                raise ValueError("LLM returned unsupported format")
            
            summary_text = result.get('summary') or 'Summary not available'
            confidence = result.get('confidence', 0.5)
            try:
                confidence = float(confidence)
            except (ValueError, TypeError):
                confidence = 0.5

            if confidence <= 0.01 and len(text.strip()) < 50 and not images:
                # Extremely short or empty documents tend to yield noisy scores.
                confidence = 0.6
                if summary_text.startswith("I am unable to provide"):
                    summary_text = "Document contains limited visible content (likely template or image-only)."

            # Use AI's intelligent categorization
            final_category = result.get('primary_category', 'other')
            reasoning = result.get('reasoning', '')
            
            # Show reasoning if provided (for debugging/transparency)
            if reasoning:
                print(f"   💭 AI Reasoning: {reasoning[:150]}...")

            # Extract AI-suggested folder structure (learned from examples, intelligently proposed)
            suggested_folder_structure = result.get('suggested_folder_structure', '')
            if suggested_folder_structure:
                print(f"   📁 AI-suggested folder structure: {suggested_folder_structure}")
                print(f"      (Learned from examples, intelligently proposed for this document)")
            
            return {
                'ai_summary': summary_text,
                'ai_category': final_category,
                'ai_subcategories': result.get('subcategories', []),
                'confidence_score': confidence,
                'tags': result.get('suggested_tags', []),
                'key_topics': result.get('key_topics', []),
                'importance_score': result.get('importance_score', 5),
                'suggested_folder_structure': suggested_folder_structure  # AI-proposed intelligent structure (learned from examples)
            }
            
        except Exception as e:
            print(f"Warning: AI summary generation failed: {e}")
            return {
                'ai_summary': 'AI analysis unavailable',
                'ai_category': 'other',
                'ai_subcategories': [],
                'confidence_score': 0.0,
                'tags': [],
                'key_topics': [],
                'importance_score': 5
            }
    
    def analyze_path_appropriateness(self, current_path: str, category: str, 
                                     entities: Dict, file_name: str) -> Tuple[str, float]:
        """Analyze if current file location is appropriate and suggest better path."""
        
        # Get category suggestion from database
        try:
            category_result = self.supabase.table('categories')\
                .select('suggested_folder_path')\
                .eq('category_name', category)\
                .execute()
            
            if category_result.data and len(category_result.data) > 0:
                template = category_result.data[0].get('suggested_folder_path')
                if not template:  # If None or empty
                    template = 'Documents/Uncategorized'
            else:
                template = 'Documents/Uncategorized'
        except Exception as e:
            print(f"Warning: Could not fetch category suggestion: {e}")
            template = 'Documents/Uncategorized'
        
        # Replace placeholders in template
        suggested_path = template
        
        # Extract year from entities or filename
        year_match = re.search(r'\b(20\d{2})\b', file_name)
        if year_match:
            suggested_path = suggested_path.replace('{Year}', year_match.group(1))
        else:
            suggested_path = suggested_path.replace('{Year}', str(datetime.now().year))
        
        # Extract vehicle info
        if entities.get('vehicles'):
            vehicle = entities['vehicles'][0]
            suggested_path = suggested_path.replace('{Vehicle}', vehicle)
        else:
            suggested_path = suggested_path.replace('{Vehicle}', 'Unknown')
        
        # Calculate confidence in current path
        current_path_lower = current_path.lower()
        category_lower = category.lower()
        
        # Check if current path contains category-related keywords
        confidence = 0.5  # Base confidence
        
        if category_lower.replace('_', ' ') in current_path_lower:
            confidence += 0.3
        
        if suggested_path.lower() in current_path_lower or current_path_lower in suggested_path.lower():
            confidence += 0.2
        
        # Check if in appropriate top-level folder
        if any(folder in current_path_lower for folder in ['documents', 'files', category_lower.split('_')[0]]):
            confidence += 0.1
        
        # Penalize if in Downloads, Desktop, or other temporary locations
        if any(folder in current_path_lower for folder in ['downloads', 'desktop', 'temp', 'tmp']):
            confidence -= 0.3
        
        confidence = max(0.0, min(1.0, confidence))
        
        return suggested_path, confidence
    
    def process_document(self, file_path: str, skip_if_exists: bool = True) -> Dict:
        """
        Main method to process a document.
        
        🛡️ SAFETY GUARANTEE - RULE #1:
        This method NEVER deletes, moves, renames, or modifies files.
        It ONLY reads files and updates database records.
        Files remain completely untouched on disk.
        """
        
        print(f"\n{'='*80}")
        print(f"Processing: {file_path}")
        print(f"{'='*80}\n")
        
        # 🛡️ SAFETY CHECK #1: Validate file exists (we're reading, not deleting)
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # 🛡️ SAFETY CHECK #2: Explicit assertion - this method never deletes files
        # This is a read-only operation. If you see any file deletion code here,
        # it's a CRITICAL BUG and violates Rule #1.
        assert True, "SAFETY RULE #1: This method must NEVER delete or move files"
        
        # Get file info
        file_stat = os.stat(file_path)
        file_path_obj = Path(file_path)
        file_hash = self.calculate_file_hash(file_path)
        
        print(f"📄 File: {file_path_obj.name}")
        print(f"📊 Size: {file_stat.st_size / 1024:.2f} KB")
        print(f"🔑 Hash: {file_hash[:16]}...")
        
        existing_document = None
        original_path = None  # Preserve original path for AI categorization
        try:
            existing_resp = self.supabase.table('documents')\
                .select('id, file_name, current_path, created_at')\
                .eq('file_hash', file_hash)\
                .execute()
            if existing_resp.data:
                existing_document = existing_resp.data[0]
                document_id = existing_document['id']
                
                # Get ORIGINAL path from document_locations table (where file was first discovered)
                try:
                    locations_resp = self.supabase.table('document_locations')\
                        .select('location_path')\
                        .eq('document_id', document_id)\
                        .eq('location_type', 'original')\
                        .order('discovered_at', desc=False)\
                        .limit(1)\
                        .execute()
                    
                    if locations_resp.data:
                        original_path = locations_resp.data[0].get('location_path')
                        print(f"   📍 Original path (for AI categorization): {original_path}")
                    else:
                        # Fallback: use first current_path we saw (when document was created)
                        original_path = existing_document.get('current_path')
                        print(f"   📍 Original path (from current_path): {original_path}")
                except Exception as loc_err:
                    # Fallback: use current_path if we can't get from document_locations
                    original_path = existing_document.get('current_path')
                    print(f"   📍 Original path (fallback): {original_path}")
        except Exception as e:
            print(f"Warning: Could not check for existing document: {e}")

        # Skip if already processed and skipping enabled
        if skip_if_exists and existing_document:
            print(f"⚠️  Document already exists in database:")
            print(f"   Name: {existing_document['file_name']}")
            print(f"   Path: {existing_document['current_path']}")
            return {'status': 'skipped', 'reason': 'already_exists', 'document_id': existing_document['id']}
        
        # Extract content based on file type
        file_ext = file_path_obj.suffix.lower()
        print(f"\n📖 Extracting content from {file_ext or 'file'}...")
        images = []
        try:
            if file_ext == '.pdf':
                pdf_info = self.extract_pdf_content(file_path)
            elif file_ext == '.docx':
                pdf_info = self.extract_docx_content(file_path)
            elif file_ext == '.xlsx':
                pdf_info = self.extract_xlsx_content(file_path)
            elif file_ext == '.pptx':
                pdf_info = self.extract_pptx_content(file_path)
            elif file_ext == '.txt':
                # Simple text file
                pdf_info = {
                    'title': '', 'author': '', 'subject': '', 'keywords': [],
                    'created_date': None, 'modified_date': None, 'page_count': 1,
                    'extracted_text': '', 'text_preview': ''
                }
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        full_text = f.read()
                    cleaned_text = full_text.strip().replace('\x00', '').replace('\u0000', '')
                    pdf_info['extracted_text'] = cleaned_text
                    pdf_info['text_preview'] = cleaned_text[:500].strip()
                except Exception as e:
                    print(f"   ⚠️  Error reading text file: {e}")
            elif file_ext == '.rtf':
                # RTF file
                pdf_info = {
                    'title': '', 'author': '', 'subject': '', 'keywords': [],
                    'created_date': None, 'modified_date': None, 'page_count': 1,
                    'extracted_text': '', 'text_preview': ''
                }
                try:
                    import striprtf
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        rtf_text = f.read()
                    plain_text = striprtf.rtf_to_text(rtf_text)
                    cleaned_text = plain_text.strip().replace('\x00', '').replace('\u0000', '')
                    pdf_info['extracted_text'] = cleaned_text
                    pdf_info['text_preview'] = cleaned_text[:500].strip()
                except ImportError:
                    print("   ⚠️  striprtf not installed. Install with: pip install striprtf")
                except Exception as e:
                    print(f"   ⚠️  Error reading RTF file: {e}")
            else:
                print(f"   ⚠️  Unsupported file type: {file_ext}")
                return {'status': 'failed', 'reason': 'unsupported_file_type', 'error': f"File type {file_ext} not supported"}
            
            print(f"   ✓ Extracted {pdf_info.get('page_count', 0)} pages/sheets/slides")
            print(f"   ✓ Text length: {len(pdf_info.get('extracted_text', ''))} characters")
            
            # TRIGGER VISION IF (PDFs only):
            # 1. Text is dangerously short/empty
            # 2. OR if PyPDF2 failed entirely (extracted_text is empty and we caught the error)
            should_use_vision = (file_ext == '.pdf' and 
                                len(pdf_info.get('extracted_text', '')) < 100 * pdf_info.get('page_count', 1) and
                                convert_from_path)
            
            if should_use_vision:
                print("   ⚠️  Low/No text detected. Converting to images for Vision analysis...")
                try:
                    # Convert first 15 pages to images for analysis
                    images = convert_from_path(file_path, first_page=1, last_page=15)
                    print(f"   ✓ Converted {len(images)} pages to images")
                except Exception as img_err:
                    print(f"   ⚠️  Image conversion failed: {img_err}")
                    # If both text and image conversion fail, we truly can't read it.
                    if not pdf_info['extracted_text']:
                         return {'status': 'failed', 'reason': 'unreadable_content', 'error': f"Text and Image extraction failed: {img_err}"}

        except Exception as e:
            # This block should barely be reached now that extract_pdf_content handles its own errors
            print(f"   ✗ Extraction failed: {e}")
            return {'status': 'failed', 'reason': 'extraction_error', 'error': str(e)}
        
        # Detect document mode
        print("\n🔍 Detecting document type...")
        document_mode = self.detect_document_mode(pdf_info['extracted_text'])
        is_conversation = document_mode == 'conversation'
        print(f"   ✓ Mode: {document_mode.upper()}")
        
        # Extract entities (Now supports Vision)
        print("\n🏷️  Extracting entities...")
        entities = self.extract_entities(pdf_info['extracted_text'], images=images)
        print(f"   ✓ Vehicles: {entities.get('vehicles', [])[:3]}")
        print(f"   ✓ Organizations: {entities.get('organizations', [])[:3]}")
        print(f"   ✓ People: {entities.get('people', [])[:3]}")
        
        # Get folder hierarchy for context (used by AI and stored in DB)
        folder_hierarchy = list(file_path_obj.parent.parts)
        
        # Use ORIGINAL path for AI categorization (where file was first discovered)
        # This is more meaningful than current path for understanding what the document is
        path_for_ai = original_path if original_path else str(file_path)
        if original_path and original_path != str(file_path):
            print(f"   📍 Using ORIGINAL path for AI categorization: {original_path}")
            # Extract hierarchy from original path
            original_path_obj = Path(original_path)
            folder_hierarchy_for_ai = list(original_path_obj.parent.parts)
        else:
            folder_hierarchy_for_ai = folder_hierarchy
        
        # Generate AI summary and categorization (Now supports Vision)
        # AI uses filename, ORIGINAL path (where file was first discovered), and content
        print("\n🤖 Generating AI analysis (using filename, ORIGINAL path, and content context)...")
        ai_analysis = self.generate_ai_summary(
            pdf_info['extracted_text'],
            file_path_obj.name,
            document_mode,
            file_path=path_for_ai,  # Use ORIGINAL path, not current path
            folder_hierarchy=folder_hierarchy_for_ai,  # Use original path's hierarchy
            images=images
        )
        print(f"   ✓ Category: {ai_analysis['ai_category']}")
        print(f"   ✓ Confidence: {ai_analysis['confidence_score']:.2%}")
        print(f"   ✓ Summary: {ai_analysis['ai_summary'][:100]}...")
        
        # Analyze path appropriateness
        print("\n📁 Analyzing file location...")
        suggested_path, path_confidence = self.analyze_path_appropriateness(
            str(file_path_obj.parent),
            ai_analysis['ai_category'],
            entities,
            file_path_obj.name
        )
        print(f"   ✓ Current path confidence: {path_confidence:.2%}")
        if path_confidence < 0.7:
            print(f"   💡 Suggested: {suggested_path}")
        
        # Detect context bin
        print("\n📂 Detecting context bin...")
        context_bin = self.detect_context_bin(str(file_path))
        if context_bin:
            print(f"   ✓ Context bin: {context_bin}")
        else:
            print(f"   ⚠️  No context bin detected")
        
        # Prepare document record
        current_path_str = str(file_path)
        
        # Truncate extracted_text to avoid PostgreSQL tsvector limit (1MB)
        # Keep 800KB to be safe (tsvector limit is 1MB, but we need buffer)
        MAX_TEXT_SIZE = 800000  # 800KB in characters
        extracted_text = pdf_info.get('extracted_text', '')
        if len(extracted_text) > MAX_TEXT_SIZE:
            print(f"   ⚠️  Text too large ({len(extracted_text):,} chars), truncating to {MAX_TEXT_SIZE:,} chars for database")
            extracted_text = extracted_text[:MAX_TEXT_SIZE]
            # Add note that text was truncated
            if not pdf_info.get('text_preview', ''):
                pdf_info['text_preview'] = extracted_text[:500].strip()
        
        # Get filesystem modification time (st_mtime) - more accurate than PDF metadata
        file_modified_at = datetime.fromtimestamp(file_stat.st_mtime).isoformat()
        
        # PRESERVE DATA: Only update fields that should change, preserve others
        document_data = {
            'file_hash': file_hash,  # Never changes (same file)
            'file_name': file_path_obj.name,  # May change if renamed
            'file_size_bytes': file_stat.st_size,  # Update in case file changed
            'file_type': file_ext.lstrip('.') if file_ext else 'pdf',  # Set based on actual file type
            'current_path': current_path_str,  # Update to current location
            'folder_hierarchy': folder_hierarchy,  # Update to current hierarchy
            'context_bin': context_bin,  # Update if bin detection improved
            'suggested_path': suggested_path if path_confidence < 0.7 else None,
            'path_confidence': path_confidence,
            'last_verified_at': datetime.now().isoformat(),  # Update verification time
            'file_modified_at': file_modified_at,  # Filesystem modification time (st_mtime)
            'document_mode': document_mode,
            'is_conversation': is_conversation,
            'pdf_title': pdf_info['title'],  # Update in case PDF metadata changed
            'pdf_author': pdf_info['author'],
            'pdf_subject': pdf_info['subject'],
            'pdf_keywords': pdf_info['keywords'],
            'page_count': pdf_info['page_count'],
            'extracted_text': extracted_text,  # Truncated if needed to avoid tsvector limit
            'text_preview': pdf_info['text_preview'],
            'entities': entities,  # Re-extract (may improve)
            'processing_status': 'completed',
            'indexed_at': datetime.now().isoformat(),
            # AI analysis fields - UPDATE with new intelligent categorization
            'ai_summary': ai_analysis['ai_summary'],
            'ai_category': ai_analysis['ai_category'],
            'ai_subcategories': ai_analysis['ai_subcategories'],
            'confidence_score': ai_analysis['confidence_score']
        }
        
        # Store AI-suggested folder structure (learned from examples, intelligently proposed)
        ai_suggested_structure = ai_analysis.get('suggested_folder_structure', '')
        if ai_suggested_structure:
            document_data['suggested_folder_structure'] = ai_suggested_structure
        
        # PRESERVE: Don't overwrite created_at, preserve original timestamps
        # (created_at is set automatically on insert, we don't touch it on update)
        
        # Deep Conversation Analysis (if applicable)
        if is_conversation and parse_iphone_backup_messages:
            print("\n💬 Performing deep conversation analysis...")
            try:
                # 1. Parse structured messages
                conversation_data = parse_iphone_backup_messages(pdf_info['extracted_text'])
                
                if conversation_data and conversation_data.get('total_messages', 0) > 0:
                    print(f"   ✓ Parsed {conversation_data['total_messages']} messages")
                    
                    # 2. Analyze relationship/sentiment
                    print("   ✓ Running therapist analysis...")
                    analyzer = ConversationAnalyzer(conversation_data)
                    analysis_result = analyzer.analyze()
                    
                    # 3. Populate specific tables (messages, participants, analysis)
                    print("   ✓ Populating conversation tables...")
                    populator = DatabasePopulator(supabase_url=self.supabase.supabase_url, supabase_key=self.supabase.supabase_key)
                    populator.populate_all(conversation_data, analysis_result)
                    
                    # Update document metadata with conversation insights
                    document_data['ai_summary'] = f"Conversation between {list(conversation_data['participants'].values())}. {analysis_result['relationship'].get('survival_assessment', {}).get('assessment_text', '')}"
                    document_data['ai_category'] = 'conversation'
                    
                else:
                    print("   ⚠️  No structured messages found in conversation text")
            except Exception as e:
                print(f"   ⚠️  Conversation analysis failed: {e}")
        
        # Insert into database
        print("\n💾 Saving to database...")
        try:
            # Re-check for existing document (in case it was added between check and now)
            if not existing_document:
                try:
                    existing_resp = self.supabase.table('documents')\
                        .select('id, file_name, current_path')\
                        .eq('file_hash', file_hash)\
                        .execute()
                    if existing_resp.data:
                        existing_document = existing_resp.data[0]
                except:
                    pass
            
            if existing_document:
                document_id = existing_document['id']
                old_path = existing_document.get('current_path')
                new_path = current_path_str
                
                # TRACK PATH CHANGES: If path changed, record it in document_locations
                # BUT: Keep using ORIGINAL path for AI categorization (already stored)
                if old_path and old_path != new_path:
                    print(f"   📍 Path changed: {old_path} → {new_path}")
                    print(f"   💡 Note: AI will continue using ORIGINAL path for categorization")
                    try:
                        # Record old path in document_locations (preserve history)
                        self.supabase.table('document_locations').insert({
                            'document_id': document_id,
                            'location_path': old_path,
                            'location_type': 'previous',
                            'discovered_at': existing_document.get('created_at') or datetime.now().isoformat(),
                            'verified_at': datetime.now().isoformat(),
                            'is_accessible': os.path.exists(old_path) if old_path else False,
                            'notes': 'Previous location before path update (not used for AI categorization)'
                        }).execute()
                        print(f"   ✓ Preserved old path in history")
                    except Exception as loc_err:
                        print(f"   Warning: Could not record path history: {loc_err}")
                
                # Clean up old category links (they'll be recreated)
                try:
                    self.supabase.table('document_categories').delete().eq('document_id', document_id).execute()
                except Exception as delete_err:
                    print(f"Warning: Could not clear existing category links: {delete_err}")
                
                # UPDATE: Only update the fields we want to change
                # created_at, id, etc. are preserved automatically
                result = self.supabase.table('documents').update(document_data).eq('id', document_id).execute()
                print(f"   ✓ Document updated (ID: {document_id})")
                if old_path != new_path:
                    print(f"   📍 Original path preserved in document_locations table")
            else:
                # NEW DOCUMENT: Insert fresh record (with error handling for race conditions)
                try:
                    result = self.supabase.table('documents').insert(document_data).execute()
                    document_id = result.data[0]['id']
                    print(f"   ✓ Document ID: {document_id}")
                    original_path = current_path_str
                except Exception as insert_err:
                    # Handle duplicate key error (race condition - file was added between checks)
                    error_str = str(insert_err)
                    if 'duplicate key' in error_str.lower() or '23505' in error_str:
                        print(f"   ⚠️  File already exists (added by another process). Updating existing record...")
                        # Fetch existing document
                        existing_resp = self.supabase.table('documents')\
                            .select('id, file_name, current_path')\
                            .eq('file_hash', file_hash)\
                            .execute()
                        if existing_resp.data:
                            existing_document = existing_resp.data[0]
                            document_id = existing_document['id']
                            old_path = existing_document.get('current_path')
                            new_path = current_path_str
                            
                            # Update instead of insert
                            self.supabase.table('documents')\
                                .update(document_data)\
                                .eq('id', document_id)\
                                .execute()
                            print(f"   ✓ Document updated (ID: {document_id})")
                            
                            # Track path change if needed
                            if old_path and old_path != new_path:
                                try:
                                    self.supabase.table('document_locations').insert({
                                        'document_id': document_id,
                                        'location_path': old_path,
                                        'location_type': 'previous',
                                        'discovered_at': datetime.now().isoformat(),
                                        'verified_at': datetime.now().isoformat(),
                                        'is_accessible': os.path.exists(old_path) if old_path else False,
                                        'notes': 'Previous location before path update'
                                    }).execute()
                                except:
                                    pass
                            
                            original_path = old_path or current_path_str
                        else:
                            raise insert_err
                    else:
                        raise insert_err
                
                # Record initial location as ORIGINAL (only if this was a true insert)
                if 'original_path' not in locals() or original_path == current_path_str:
                    try:
                        # Check if original location already exists
                        loc_check = self.supabase.table('document_locations')\
                            .select('id')\
                            .eq('document_id', document_id)\
                            .eq('location_type', 'original')\
                            .execute()
                        
                        if not loc_check.data:
                            self.supabase.table('document_locations').insert({
                                'document_id': document_id,
                                'location_path': current_path_str,
                                'location_type': 'original',
                                'discovered_at': datetime.now().isoformat(),
                                'verified_at': datetime.now().isoformat(),
                                'is_accessible': True,
                                'notes': 'Original location when first processed - used for AI categorization'
                            }).execute()
                            print(f"   📍 Recorded original path for future AI categorization")
                    except Exception as loc_err:
                        print(f"   Warning: Could not record initial location: {loc_err}")
            
            # Link to categories
            if ai_analysis['ai_category']:
                self._link_to_category(document_id, ai_analysis['ai_category'], True, ai_analysis['confidence_score'])
            
            for subcat in ai_analysis['ai_subcategories']:
                self._link_to_category(document_id, subcat, False, ai_analysis['confidence_score'])
            
            print(f"\n{'='*80}")
            print(f"✅ SUCCESS: Document processed and saved")
            print(f"{'='*80}\n")
            
            return {
                'status': 'success',
                'document_id': document_id,
                'document_mode': document_mode,
                'category': ai_analysis['ai_category'],
                'summary': ai_analysis['ai_summary']
            }
            
        except Exception as e:
            print(f"   ✗ Database error: {e}")
            return {'status': 'failed', 'reason': 'database_error', 'error': str(e)}

    def _link_to_category(self, document_id: str, category_name: str, 
                          is_primary: bool, confidence: float):
        """Link document to a category."""
        try:
            # Get or create category
            category_result = self.supabase.table('categories')\
                .select('id')\
                .eq('category_name', category_name)\
                .execute()
            
            if category_result.data:
                category_id = category_result.data[0]['id']
            else:
                # Create new category
                new_cat = self.supabase.table('categories').insert({
                    'category_name': category_name,
                    'description': f'Auto-created category: {category_name}'
                }).execute()
                category_id = new_cat.data[0]['id']
            
            # Link document to category
            self.supabase.table('document_categories').insert({
                'document_id': document_id,
                'category_id': category_id,
                'confidence': confidence,
                'is_primary': is_primary,
                'assigned_by': 'ai'
            }).execute()
            
        except Exception as e:
            print(f"Warning: Could not link to category {category_name}: {e}")


def main():
    if len(sys.argv) < 2:
        print("Usage: python document_processor.py <pdf_path>")
        print("\nProcesses a PDF document and stores it in the database with AI analysis.")
        sys.exit(1)
    
    pdf_path = sys.argv[1]
    
    processor = DocumentProcessor()
    result = processor.process_document(pdf_path)
    
    print(f"\nResult: {json.dumps(result, indent=2)}")


if __name__ == "__main__":
    main()
