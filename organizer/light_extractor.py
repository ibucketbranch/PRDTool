"""Lightweight file signal collector for Goldilocks ingestion.

Extracts three classification signals per file:
1. Filename (parsed, split on delimiters)
2. Path context (parent folder names)
3. Content preview (first 500 chars from text-accessible files)

No OCR. No full-document parsing. Images/video/audio classified by
filename + path only.
"""

import os
import subprocess
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

MAX_PREVIEW_CHARS = 500
MAX_FILE_SIZE_BYTES = 5 * 1024 * 1024  # 5MB cap for content reads

TEXT_EXTENSIONS = frozenset({
    ".txt", ".md", ".csv", ".log", ".rtf",
    ".json", ".xml", ".yaml", ".yml",
    ".html", ".htm", ".py", ".js", ".ts",
    ".sh", ".cfg", ".ini", ".conf",
})

TEXTUTIL_EXTENSIONS = frozenset({
    ".doc", ".key", ".numbers", ".pages",
})

SKIP_CONTENT_EXTENSIONS = frozenset({
    ".png", ".jpg", ".jpeg", ".gif", ".tiff", ".bmp", ".heic", ".webp",
    ".mp4", ".mov", ".m4v", ".avi", ".mkv",
    ".mp3", ".m4a", ".wav", ".aac", ".flac",
    ".zip", ".gz", ".tar", ".dmg", ".rar", ".7z",
    ".exe", ".app", ".pkg", ".iso", ".bin",
    ".icloud",
})


@dataclass
class FileSignals:
    file_path: str
    filename: str
    parent_folders: list[str] = field(default_factory=list)
    extension: str = ""
    file_size: int = 0
    modified_date: str = ""
    content_preview: str = ""
    extraction_method: str = "filename_only"


def collect_signals(file_path: str) -> FileSignals:
    """Collect classification signals from a file.

    Args:
        file_path: Absolute path to the file.

    Returns:
        FileSignals with filename, path context, and content preview.
    """
    p = Path(file_path)
    ext = p.suffix.lower()

    stat = p.stat()
    modified = datetime.fromtimestamp(stat.st_mtime).isoformat()

    icloud_root = Path.home() / "Library/Mobile Documents/com~apple~CloudDocs"
    try:
        rel = p.relative_to(icloud_root)
        parents = [part for part in rel.parent.parts if part]
    except ValueError:
        parents = [part for part in p.parent.parts[-3:] if part]

    signals = FileSignals(
        file_path=str(p),
        filename=p.name,
        parent_folders=parents,
        extension=ext,
        file_size=stat.st_size,
        modified_date=modified,
    )

    if ext in SKIP_CONTENT_EXTENSIONS or stat.st_size > MAX_FILE_SIZE_BYTES:
        signals.extraction_method = "filename_only"
        return signals

    preview = _extract_preview(p, ext)
    if preview:
        signals.content_preview = preview[:MAX_PREVIEW_CHARS]

    return signals


def _extract_preview(path: Path, ext: str) -> str:
    """Extract content preview based on file type."""
    try:
        if ext == ".pdf":
            return _extract_pdf(path)
        elif ext == ".docx":
            return _extract_docx(path)
        elif ext == ".xlsx":
            return _extract_xlsx(path)
        elif ext == ".pptx":
            return _extract_pptx(path)
        elif ext in TEXTUTIL_EXTENSIONS:
            return _extract_textutil(path)
        elif ext in TEXT_EXTENSIONS:
            return _extract_plain(path)
        else:
            return _extract_plain(path)
    except Exception:
        return ""


def _extract_pdf(path: Path) -> str:
    import pdfplumber

    with pdfplumber.open(path) as pdf:
        if not pdf.pages:
            return ""
        text = pdf.pages[0].extract_text() or ""
        return text.strip()


def _extract_docx(path: Path) -> str:
    import docx

    doc = docx.Document(path)
    parts: list[str] = []
    for para in doc.paragraphs[:5]:
        if para.text.strip():
            parts.append(para.text.strip())
        if len("\n".join(parts)) >= MAX_PREVIEW_CHARS:
            break
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    parts.append(f"Sheets: {', '.join(wb.sheetnames)}")
    ws = wb.active
    if ws:
        for row in ws.iter_rows(max_row=3, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides[:2]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if len("\n".join(parts)) >= MAX_PREVIEW_CHARS:
                break
    return "\n".join(parts)


def _extract_textutil(path: Path) -> str:
    """Use macOS built-in textutil to extract text from .doc/.key/.numbers/.pages."""
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, text=True, timeout=10,
        )
        return result.stdout.strip() if result.returncode == 0 else ""
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return ""


def _extract_plain(path: Path) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(MAX_PREVIEW_CHARS).strip()
    except (OSError, UnicodeDecodeError):
        return ""
